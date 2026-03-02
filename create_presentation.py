#!/usr/bin/env python3
"""Generate SC-Net CAD Diagnosis progress update PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Color scheme
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF4, 0xF7)
ACCENT_BLUE = RGBColor(0x3A, 0x7C, 0xBD)
MEDIUM_GRAY = RGBColor(0x6B, 0x7B, 0x8D)
TABLE_HEADER_BG = RGBColor(0x1B, 0x3A, 0x5C)
TABLE_ALT_BG = RGBColor(0xE8, 0xEE, 0xF4)
GREEN = RGBColor(0x27, 0xAE, 0x60)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, font_name="Calibri", spacing=Pt(6)):
    """Add a bulleted list text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        p.level = 0
        # Bullet character
        p.text = f"\u2022  {item}"
    return txBox


def create_title_slide(prs):
    """Slide 1: Title + Project Overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    set_slide_bg(slide, DARK_BLUE)

    # Accent bar at top
    add_shape_bg(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT_BLUE)

    # Title
    add_textbox(slide, Inches(1), Inches(1.0), Inches(11), Inches(1),
                "SC-Net: Spatio-Temporal Contrast Network\nfor CAD Diagnosis",
                font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

    # Subtitle
    add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.6),
                "MICCAI 2024 Paper Implementation  \u2014  Progress Update",
                font_size=20, color=ACCENT_BLUE, bold=False)

    # Divider line
    add_shape_bg(slide, Inches(1), Inches(3.3), Inches(4), Inches(0.03), ACCENT_BLUE)

    # Brief description
    add_textbox(slide, Inches(1), Inches(3.6), Inches(10.5), Inches(0.8),
                "Dual-branch DETR-style architecture for coronary artery disease diagnosis from CCTA",
                font_size=17, color=MEDIUM_GRAY, bold=False)

    # Three core components - cards
    components = [
        ("Clinically-Credible\nData Augmentation",
         "Lesion recombination on CPR\nvolumes for data-scarce settings"),
        ("Spatio-Temporal\nDual-Task Learning",
         "Object detection (spatial) +\nsampling-point classification (temporal)"),
        ("Dual-Task Prediction-\nContrast Optimization",
         "Cross-supervision via detached\npseudo-labels between branches"),
    ]

    card_width = Inches(3.3)
    card_height = Inches(2.2)
    start_x = Inches(1)
    card_y = Inches(4.5)
    gap = Inches(0.4)

    for i, (title, desc) in enumerate(components):
        x = start_x + i * (card_width + gap)
        # Card background
        card = add_shape_bg(slide, x, card_y, card_width, card_height,
                            RGBColor(0x24, 0x4E, 0x78))
        card.line.color.rgb = ACCENT_BLUE
        card.line.width = Pt(1)

        # Component number
        add_textbox(slide, x + Inches(0.2), card_y + Inches(0.15),
                    Inches(0.5), Inches(0.4),
                    str(i + 1), font_size=24, color=ACCENT_BLUE, bold=True)

        # Component title
        add_textbox(slide, x + Inches(0.2), card_y + Inches(0.5),
                    card_width - Inches(0.4), Inches(0.8),
                    title, font_size=16, color=WHITE, bold=True)

        # Component description
        add_textbox(slide, x + Inches(0.2), card_y + Inches(1.35),
                    card_width - Inches(0.4), Inches(0.8),
                    desc, font_size=12, color=MEDIUM_GRAY)


def create_bugfix_slide(prs):
    """Slide 2: Implementation Progress & Bug Fixes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BLUE)

    # Accent bar
    add_shape_bg(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT_BLUE)

    # Title
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Implementation Progress & Bug Fixes",
                font_size=30, color=WHITE, bold=True)

    add_shape_bg(slide, Inches(0.8), Inches(1.05), Inches(3), Inches(0.03), ACCENT_BLUE)

    # Left column: Bug fixes
    col1_x = Inches(0.8)
    col1_w = Inches(5.5)

    add_textbox(slide, col1_x, Inches(1.3), col1_w, Inches(0.5),
                "13+ Architecture Bugs Identified & Fixed",
                font_size=18, color=ACCENT_BLUE, bold=True)

    fixes = [
        "nn.ModuleList for extraction blocks (weights were invisible to optimizer)",
        "Learnable parameters: feature weights, view fusion as nn.Parameter",
        "Fixed query embeddings via nn.Embedding (was random every forward pass)",
        "Conv3d spatial flattening projection (defined but never called)",
        "Gradient detachment in contrastive loss (circular gradient flow)",
        "Box format consistency: center-width throughout pipeline",
        "Device handling: targets moved to output device in loss functions",
        "Deep copy targets before each loss term (in-place mutation bug)",
    ]
    add_bullet_list(slide, col1_x, Inches(1.8), col1_w, Inches(3.5),
                    fixes, font_size=13, color=RGBColor(0xCC, 0xD5, 0xE0), spacing=Pt(4))

    # Right column: Key discovery + training pipeline
    col2_x = Inches(7.0)
    col2_w = Inches(5.5)

    # Discovery box
    disc_y = Inches(1.3)
    disc_h = Inches(2.4)
    disc_bg = add_shape_bg(slide, col2_x, disc_y, col2_w, disc_h,
                           RGBColor(0x24, 0x4E, 0x78))
    disc_bg.line.color.rgb = ORANGE
    disc_bg.line.width = Pt(1.5)

    add_textbox(slide, col2_x + Inches(0.2), disc_y + Inches(0.1),
                col2_w - Inches(0.4), Inches(0.5),
                "Key Discovery: Paper Code \u2260 Paper Equations",
                font_size=16, color=ORANGE, bold=True)

    disc_items = [
        "Paper equations: \u03bb_L1 = 5, \u03bb_IoU = 2",
        "Paper code actually uses: 1:1:1 equal weighting",
        "Our \"fixes\" matching equations broke convergence",
        "Resolution: Reverted to paper code weights (1:1:1)",
    ]
    add_bullet_list(slide, col2_x + Inches(0.2), disc_y + Inches(0.55),
                    col2_w - Inches(0.4), Inches(1.6),
                    disc_items, font_size=13, color=RGBColor(0xCC, 0xD5, 0xE0), spacing=Pt(3))

    # Training pipeline box
    pipe_y = Inches(4.0)
    pipe_h = Inches(3.0)
    pipe_bg = add_shape_bg(slide, col2_x, pipe_y, col2_w, pipe_h,
                           RGBColor(0x24, 0x4E, 0x78))
    pipe_bg.line.color.rgb = ACCENT_BLUE
    pipe_bg.line.width = Pt(1)

    add_textbox(slide, col2_x + Inches(0.2), pipe_y + Inches(0.1),
                col2_w - Inches(0.4), Inches(0.5),
                "Two-Stage Training Pipeline",
                font_size=16, color=ACCENT_BLUE, bold=True)

    pipe_items = [
        "Stage 1 \u2014 Pre-train on augmented data (3-class plaque)",
        "Stage 2 \u2014 Fine-tune on clinical data (6-class: stenosis \u00d7 plaque)",
        "Head reinitialization (3\u21926 classes) is by design in paper code",
        "AdamW optimizer, CosineAnnealingLR, gradient clipping = 0.1",
    ]
    add_bullet_list(slide, col2_x + Inches(0.2), pipe_y + Inches(0.55),
                    col2_w - Inches(0.4), Inches(2.0),
                    pipe_items, font_size=13, color=RGBColor(0xCC, 0xD5, 0xE0), spacing=Pt(3))

    # Bottom left: Phase summary
    phase_y = Inches(5.5)
    add_textbox(slide, col1_x, phase_y, col1_w, Inches(0.4),
                "Development Phases",
                font_size=16, color=ACCENT_BLUE, bold=True)

    phases = [
        "Phase 1: Initial codebase setup (Jan 2025)",
        "Phase 2: 15+ critical bug fixes across 6 files (Jan 2025 \u2013 Feb 2026)",
        "Phase 3: Training infrastructure, DDP, evaluation pipeline (Feb 2026)",
        "Phase 4: Root cause analysis \u2014 paper code vs. paper equations (Feb 2026)",
    ]
    add_bullet_list(slide, col1_x, phase_y + Inches(0.35), col1_w, Inches(1.5),
                    phases, font_size=13, color=RGBColor(0xCC, 0xD5, 0xE0), spacing=Pt(3))


def create_results_slide(prs):
    """Slide 3: Current Results with table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BLUE)

    # Accent bar
    add_shape_bg(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT_BLUE)

    # Title
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Current Results",
                font_size=30, color=WHITE, bold=True)
    add_shape_bg(slide, Inches(0.8), Inches(1.05), Inches(3), Inches(0.03), ACCENT_BLUE)

    # Results table — best results from each evaluation
    rows, cols = 6, 6
    table_left = Inches(0.8)
    table_top = Inches(1.4)
    table_width = Inches(11.5)
    table_height = Inches(2.8)

    table_shape = slide.shapes.add_table(rows, cols, table_left, table_top,
                                         table_width, table_height)
    table = table_shape.table

    # Column widths
    col_widths = [Inches(2.6), Inches(1.6), Inches(1.6), Inches(1.6), Inches(1.6), Inches(2.5)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Headers
    headers = ["Run", "Stenosis ACC", "Macro F1", "Sig. F1", "Macro AUC", "Notes"]
    data = [
        ["v2-ft (baseline)", "0.316", "0.160", "0.000", "0.573", "Majority class only"],
        ["v6-ft (argmax)", "0.369", "0.288", "0.000", "0.645", "Zero Sig. predictions"],
        ["v6-ft (calibrated)", "0.470", "0.417", "0.621", "0.645", "Best result (val split)"],
        ["v6-ft (held-out, cal.)", "0.435", "0.393", "0.525", "0.604", "Unseen AP-NUH patients"],
        ["Paper target", "0.914", "\u2014", "\u2014", "\u2014", "100% data (218 patients)"],
    ]

    def style_cell(cell, text, is_header=False, highlight=False):
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        if is_header:
            p.font.bold = True
            p.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_HEADER_BG
        elif highlight:
            p.font.color.rgb = WHITE
            p.font.bold = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x24, 0x4E, 0x78)
        else:
            p.font.color.rgb = RGBColor(0xDD, 0xE4, 0xED)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x1E, 0x42, 0x66)

    for j, h in enumerate(headers):
        style_cell(table.cell(0, j), h, is_header=True)

    for i, row_data in enumerate(data):
        is_highlight = (i == 2 or i == 3 or i == 4)  # calibrated rows + paper target
        for j, val in enumerate(row_data):
            style_cell(table.cell(i + 1, j), val, highlight=is_highlight)

    # Key findings text box (left side)
    finding_y = Inches(4.5)
    finding_h = Inches(2.7)
    finding_bg = add_shape_bg(slide, Inches(0.8), finding_y, Inches(5.0), finding_h,
                              RGBColor(0x24, 0x4E, 0x78))
    finding_bg.line.color.rgb = GREEN
    finding_bg.line.width = Pt(1.5)

    add_textbox(slide, Inches(1.0), finding_y + Inches(0.1),
                Inches(4.6), Inches(0.5),
                "Calibration Impact (Held-Out Test)",
                font_size=16, color=GREEN, bold=True)

    findings = [
        "Significant AUC = 0.707 \u2014 model discriminates internally",
        "Argmax: zero Significant predictions (boundary misaligned)",
        "Thresholds: H=3.0, NS=1.0, Sig=0.346",
        "Significant F1: 0.000 \u2192 0.525 (no retraining)",
        "Macro F1: 0.210 \u2192 0.393 (+87%)",
        "Plaque AUC \u2248 0.55 (near chance) \u2014 needs training fixes",
    ]
    add_bullet_list(slide, Inches(1.0), finding_y + Inches(0.55),
                    Inches(4.6), Inches(2.0),
                    findings, font_size=12, color=RGBColor(0xCC, 0xD5, 0xE0), spacing=Pt(2))

    # ROC curve image (right side)
    import os
    roc_path = os.path.join(os.path.dirname(__file__),
                            "plots_v6ft_heldout_cal", "roc_stenosis.png")
    if os.path.exists(roc_path):
        slide.shapes.add_picture(roc_path, Inches(6.2), finding_y,
                                 Inches(3.4), Inches(2.7))

    # Confusion matrix image (far right)
    cm_path = os.path.join(os.path.dirname(__file__),
                           "plots_v6ft_heldout_cal", "confusion_stenosis.png")
    if os.path.exists(cm_path):
        slide.shapes.add_picture(cm_path, Inches(9.8), finding_y,
                                 Inches(3.2), Inches(2.7))


def create_roadmap_slide(prs):
    """Slide 4: Next Steps & Roadmap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BLUE)

    # Accent bar
    add_shape_bg(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT_BLUE)

    # Title
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Next Steps & Roadmap",
                font_size=30, color=WHITE, bold=True)
    add_shape_bg(slide, Inches(0.8), Inches(1.05), Inches(3), Inches(0.03), ACCENT_BLUE)

    # Three columns for immediate / short-term / medium-term
    columns = [
        {
            "title": "Completed",
            "color": GREEN,
            "items": [
                "Post-hoc threshold calibration\n\u2714 Sig. F1: 0.000 \u2192 0.525",
                "Held-out test evaluation on\n665 separate AP-NUH patients",
                "Root cause analysis: paper code\n\u2260 paper equations (reverted)",
                "13+ architecture bug fixes\nacross 6 core files",
            ],
        },
        {
            "title": "Next: v7 Training",
            "color": ACCENT_BLUE,
            "items": [
                "Delayed L_dc ramp schedule\n(let branches stabilize first)",
                "Confidence-gated pseudo-labels\nfor contrastive loss",
                "Class-balanced sampling to\naddress stenosis imbalance",
                "Expected: better convergence,\nhigher AUC across classes",
            ],
        },
        {
            "title": "Medium-Term",
            "color": ORANGE,
            "items": [
                "Plaque branch investigation\n(AUC \u2248 0.55, near chance)",
                "Architectural improvements:\nlarger backbone, more queries",
                "Data scaling: additional\npatient cohorts if available",
                "Full ablation study matching\npaper's Table 2 format",
            ],
        },
    ]

    col_width = Inches(3.6)
    col_height = Inches(4.5)
    start_x = Inches(0.8)
    col_y = Inches(1.4)
    gap = Inches(0.35)

    for i, col in enumerate(columns):
        x = start_x + i * (col_width + gap)

        # Card background
        card = add_shape_bg(slide, x, col_y, col_width, col_height,
                            RGBColor(0x24, 0x4E, 0x78))
        card.line.color.rgb = col["color"]
        card.line.width = Pt(1.5)

        # Title bar
        add_shape_bg(slide, x, col_y, col_width, Inches(0.5), col["color"])
        add_textbox(slide, x + Inches(0.15), col_y + Inches(0.05),
                    col_width - Inches(0.3), Inches(0.4),
                    col["title"], font_size=16, color=WHITE, bold=True,
                    alignment=PP_ALIGN.CENTER)

        # Items
        add_bullet_list(slide, x + Inches(0.2), col_y + Inches(0.6),
                        col_width - Inches(0.4), col_height - Inches(0.8),
                        col["items"], font_size=13,
                        color=RGBColor(0xCC, 0xD5, 0xE0), spacing=Pt(6))

    # Paper reference targets at bottom
    ref_y = Inches(6.2)
    ref_bg = add_shape_bg(slide, Inches(0.8), ref_y, Inches(11.5), Inches(0.9),
                          RGBColor(0x15, 0x2E, 0x4A))
    ref_bg.line.color.rgb = MEDIUM_GRAY
    ref_bg.line.width = Pt(0.5)

    add_textbox(slide, Inches(1.0), ref_y + Inches(0.1),
                Inches(11), Inches(0.3),
                "Paper Target Metrics (100% data, 218 patients)",
                font_size=14, color=ACCENT_BLUE, bold=True)

    add_textbox(slide, Inches(1.0), ref_y + Inches(0.45),
                Inches(11), Inches(0.35),
                "Stenosis:  ACC 0.928  |  F1 0.944  |  Precision 0.948  |  Recall 0.940  |  Spec 0.917          "
                "Plaque:  ACC 0.805  |  F1 0.840",
                font_size=13, color=RGBColor(0xAA, 0xBB, 0xCC))


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    create_title_slide(prs)
    create_bugfix_slide(prs)
    create_results_slide(prs)
    create_roadmap_slide(prs)

    output_path = "/home/reet/development/CAD_diagnosis/SC_Net_Progress_Update.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")


if __name__ == "__main__":
    main()
